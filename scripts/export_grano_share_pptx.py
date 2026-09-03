#!/usr/bin/env python3
"""HTMLシェアストーリーを編集可能な 16:9 PowerPoint にする（v47）。

- 背景: 写真＋ベール＋装飾のみ（文字非表示の JPEG）
- 文字: python-pptx テキストボックス（座標・スタイルは Playwright で取得）
- 1 HTML スライド = 1 PPTX スライド（frag なし）
"""
from __future__ import annotations

import http.server
import re
import shutil
import threading
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt
from playwright.sync_api import sync_playwright

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
SHOTS = Path("/tmp/grano_share_pptx_shots")
W, H = 1920, 1080
SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
SLIDE_W = Emu(int(914400 * SLIDE_W_IN))
SLIDE_H = Emu(int(914400 * SLIDE_H_IN))
EXPORT_VERSION = "v47"

HIDE_UI_CSS = """
.hint,.controls,.alt-link,.stamp,.progress{display:none!important}
.kenburns,.dust i,.line{animation:none!important}
.kenburns{transform:none!important}
.slide{opacity:0!important;visibility:hidden!important;transition:none!important;z-index:0}
.slide.active{opacity:1!important;visibility:visible!important;z-index:2}
.frag{opacity:1!important;transform:none!important;transition:none!important}
"""

HIDE_TEXT_CSS = """
.slide.active h1,.slide.active h2,.slide.active h3,.slide.active h4,
.slide.active p,.slide.active .meet-name,.slide.active .num,.slide.active .lbl,
.slide.active .corp-name,.slide.active .novas-mark,.slide.active .val h4,
.slide.active .eyebrow,.slide.active .turn,.slide.active .quote,.slide.active .lead,
.slide.active .subnames,.slide.active .catch,.slide.active .cta-line
{visibility:hidden!important}
"""

EXTRACT_TEXT_JS = """
() => {
  const slide = document.querySelector('.slide.active');
  if (!slide) return [];
  const sel = [
    '.eyebrow', '.turn', 'h1', 'h2', 'h3', 'h4',
    '.lead', '.quote', '.subnames', '.catch', '.cta-line',
    '.meet-name', '.corp-name', '.novas-mark',
    '.num', '.lbl', '.val h4'
  ].join(',');
  const seen = new Set();
  const out = [];
  for (const el of slide.querySelectorAll(sel)) {
    if (seen.has(el)) continue;
    seen.add(el);
    const r = el.getBoundingClientRect();
    if (r.width < 1 || r.height < 1) continue;
    const cs = getComputedStyle(el);
    if (cs.visibility === 'hidden' || cs.display === 'none' || cs.opacity === '0') continue;
    let text = '';
    const walk = (node) => {
      if (node.nodeType === Node.TEXT_NODE) {
        text += node.textContent;
      } else if (node.nodeName === 'BR') {
        text += '\\n';
      } else if (node.nodeType === Node.ELEMENT_NODE) {
        for (const c of node.childNodes) walk(c);
        if (['P','H1','H2','H3','H4','DIV'].includes(node.nodeName) && node !== el) {
          /* inline only */
        }
      }
    };
    walk(el);
    text = text.replace(/[ \\t]+/g, ' ').replace(/ *\\n */g, '\\n').trim();
    if (!text) continue;
    const ff = cs.fontFamily.toLowerCase();
    let family = 'Yu Gothic UI';
    if (ff.includes('mincho') || ff.includes('shippori')) family = 'Yu Mincho';
    else if (ff.includes('meiryo')) family = 'Meiryo';
    const fw = parseInt(cs.fontWeight, 10) || 400;
    const lh = parseFloat(cs.lineHeight);
    const fs = parseFloat(cs.fontSize);
    out.push({
      x: r.x, y: r.y, w: r.width, h: r.height,
      text,
      fontSize: fs,
      lineHeight: Number.isFinite(lh) ? lh : fs * 1.2,
      fontFamily: family,
      fontWeight: fw,
      color: cs.color,
      textAlign: cs.textAlign,
      letterSpacing: cs.letterSpacing,
    });
  }
  return out;
}
"""

OVERFLOW_JS = """
() => {
  const slide = document.querySelector('.slide.active');
  if (!slide) return [];
  const bad = [];
  const vw = window.innerWidth, vh = window.innerHeight;
  const sel = 'h1,h2,h3,h4,p,.meet-name,.num,.lbl,.corp-name,.novas-mark,.val h4,.eyebrow,.turn,.quote,.lead,.subnames';
  for (const el of slide.querySelectorAll(sel)) {
    const r = el.getBoundingClientRect();
    if (r.width < 1) continue;
    if (r.bottom > vh + 2 || r.right > vw + 2 || r.top < -2 || r.left < -2) {
      bad.push({tag: el.tagName, cls: el.className, top: r.top, bottom: r.bottom, right: r.right});
    }
    if (el.scrollHeight > el.clientHeight + 2) {
      bad.push({tag: el.tagName, cls: el.className, overflow: el.scrollHeight - el.clientHeight});
    }
  }
  return bad;
}
"""


def serve(directory: Path, port: int):
    class H(http.server.SimpleHTTPRequestHandler):
        def __init__(self, *a, **k):
            super().__init__(*a, directory=str(directory), **k)

        def log_message(self, *a):
            pass

    httpd = http.server.ThreadingHTTPServer(("127.0.0.1", port), H)
    t = threading.Thread(target=httpd.serve_forever, daemon=True)
    t.start()
    return httpd


def px_to_emu_x(px: float) -> int:
    return int(px / W * SLIDE_W)


def px_to_emu_y(px: float) -> int:
    return int(px / H * SLIDE_H)


def px_to_pt(px: float) -> float:
    return px * SLIDE_H_IN / H * 72


def parse_color(css_color: str) -> RGBColor:
    m = re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+)", css_color)
    if m:
        return RGBColor(int(m[1]), int(m[2]), int(m[3]))
    return RGBColor(244, 238, 228)


def align_map(text_align: str) -> PP_ALIGN:
    return {
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "end": PP_ALIGN.RIGHT,
    }.get(text_align, PP_ALIGN.LEFT)


def add_textbox(slide, item: dict):
    left = px_to_emu_x(item["x"])
    top = px_to_emu_y(item["y"])
    width = max(px_to_emu_x(item["w"]), Emu(91440))
    height = max(px_to_emu_y(item["h"]), Emu(91440))
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    lines = item["text"].split("\n")
    size_pt = px_to_pt(item["fontSize"])
    bold = item["fontWeight"] >= 600
    color = parse_color(item["color"])
    family = item["fontFamily"]
    align = align_map(item.get("textAlign", "left"))

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_before = p.space_after = Pt(0)
        if item.get("lineHeight") and item["fontSize"]:
            ratio = item["lineHeight"] / item["fontSize"]
            p.line_spacing = ratio
        f = p.font
        f.size = Pt(max(size_pt, 24))
        f.bold = bold
        f.name = family
        f.color.rgb = color


def prepare(page):
    page.evaluate("() => document.fonts.ready")
    page.evaluate(
        """() => Promise.all([...document.images].map(img => {
          if (img.complete) return Promise.resolve();
          return new Promise(r => { img.onload = r; img.onerror = r; });
        }))"""
    )
    page.wait_for_timeout(400)


def set_hide_text(page, on: bool):
    page.evaluate(
        """(active) => {
          let s = document.getElementById('pptx-hide-text');
          if (!s) {
            s = document.createElement('style');
            s.id = 'pptx-hide-text';
            document.head.appendChild(s);
          }
          s.textContent = active ? `.slide.active h1,.slide.active h2,.slide.active h3,.slide.active h4,
.slide.active p,.slide.active .meet-name,.slide.active .num,.slide.active .lbl,
.slide.active .corp-name,.slide.active .novas-mark,.slide.active .val h4,
.slide.active .eyebrow,.slide.active .turn,.slide.active .quote,.slide.active .lead,
.slide.active .subnames,.slide.active .catch,.slide.active .cta-line
{visibility:hidden!important}` : '';
        }""",
        on,
    )
    page.wait_for_timeout(120)


def show_slide(page, index: int):
    page.evaluate(
        """(i) => {
          const slides = [...document.querySelectorAll('.slide')];
          slides.forEach((s, j) => {
            s.classList.toggle('active', j === i);
            s.querySelectorAll('.frag').forEach(f => f.classList.add('on'));
          });
        }""",
        index,
    )
    page.wait_for_timeout(200)


def capture_editable(html_name: str, out_pptx: Path, port: int) -> tuple[int, int]:
    if SHOTS.exists():
        shutil.rmtree(SHOTS)
    SHOTS.mkdir(parents=True)

    textbox_total = 0
    with sync_playwright() as p:
        browser = p.chromium.launch(args=["--disable-web-security"])
        page = browser.new_page(
            viewport={"width": W, "height": H},
            device_scale_factor=1,
        )
        page.goto(f"http://127.0.0.1:{port}/{html_name}", wait_until="networkidle")
        page.add_style_tag(content=HIDE_UI_CSS)
        prepare(page)

        n = page.evaluate("() => document.querySelectorAll('.slide').length")
        slide_data: list[tuple[Path, list[dict]]] = []

        for i in range(n):
            show_slide(page, i)
            set_hide_text(page, False)
            texts = page.evaluate(EXTRACT_TEXT_JS)
            set_hide_text(page, True)
            bg_path = SHOTS / f"{i:02d}_bg.jpg"
            page.screenshot(path=str(bg_path), type="jpeg", quality=90)
            slide_data.append((bg_path, texts))
            textbox_total += len(texts)

        browser.close()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    for bg_path, texts in slide_data:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(str(bg_path), 0, 0, SLIDE_W, SLIDE_H)
        for item in texts:
            add_textbox(s, item)
    prs.save(out_pptx)
    return n, textbox_total


def check_overflow(html_name: str, port: int) -> list[dict]:
    issues: list[dict] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": W, "height": H})
        page.goto(f"http://127.0.0.1:{port}/{html_name}", wait_until="networkidle")
        page.add_style_tag(content=HIDE_UI_CSS)
        prepare(page)
        n = page.evaluate("() => document.querySelectorAll('.slide').length")
        for i in range(n):
            show_slide(page, i)
            bad = page.evaluate(OVERFLOW_JS)
            for b in bad:
                b["slide"] = i + 1
                issues.append(b)
        browser.close()
    return issues


def count_visible_chars(html_name: str, port: int) -> int:
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": W, "height": H})
        page.goto(f"http://127.0.0.1:{port}/{html_name}", wait_until="networkidle")
        page.add_style_tag(content=HIDE_UI_CSS)
        prepare(page)
        n = page.evaluate("() => document.querySelectorAll('.slide').length")
        total = 0
        for i in range(n):
            show_slide(page, i)
            t = page.evaluate(
                """() => {
                  const slide = document.querySelector('.slide.active');
                  const sel = 'h1,h2,h3,h4,p,.meet-name,.num,.lbl,.corp-name,.novas-mark,.val h4,.eyebrow,.turn,.quote,.lead,.subnames';
                  let s = '';
                  for (const el of slide.querySelectorAll(sel)) {
                    if (getComputedStyle(el).visibility === 'hidden') continue;
                    s += el.innerText;
                  }
                  return s.replace(/\\s+/g, '');
                }"""
            )
            total += len(t)
        browser.close()
    return total


def main():
    port = 18765
    httpd = serve(DOCS, port)
    try:
        issues = check_overflow("grano_share_10min.html", port)
        if issues:
            print(f"溢れ警告 {len(issues)}件:")
            for it in issues[:8]:
                print(f"  スライド{it.get('slide')}: {it}")
        else:
            print("溢れチェック: OK")

        chars = count_visible_chars("grano_share_10min.html", port)
        print(f"10分版 表示文字数（空白除く）: {chars}")

        n10, tb10 = capture_editable(
            "grano_share_10min.html", DOCS / "grano_share_10min.pptx", port
        )
        n5, tb5 = capture_editable(
            "grano_share_5min.html", DOCS / "grano_share_5min.pptx", port
        )
    finally:
        httpd.shutdown()

    print(f"export {EXPORT_VERSION}")
    print(f"10min: {n10}枚, テキストボックス {tb10}個 → {DOCS / 'grano_share_10min.pptx'}")
    print(f"5min:  {n5}枚, テキストボックス {tb5}個 → {DOCS / 'grano_share_5min.pptx'}")


if __name__ == "__main__":
    main()
